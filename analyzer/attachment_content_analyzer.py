"""
Attachment content analysis module for phishalyzer.
Extracts and analyzes text content from various file types for phishing indicators.
"""

import re
import io
import os
import zipfile
import xml.etree.ElementTree as ET
from difflib import SequenceMatcher
from collections import defaultdict

# Import compatible output system
try:
    from .compatible_output import output, print_status
    COMPATIBLE_OUTPUT = True
except ImportError:
    COMPATIBLE_OUTPUT = False

# Import existing URL analysis functionality
try:
    from . import url_extractor
    URL_ANALYSIS_AVAILABLE = True
except ImportError:
    URL_ANALYSIS_AVAILABLE = False

# Import defanging functionality
try:
    from . import defanger
    DEFANGER_AVAILABLE = True
except ImportError:
    DEFANGER_AVAILABLE = False

# Optional imports with fallbacks
try:
    import fitz  # PyMuPDF for PDF analysis
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract  # OCR for images
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    import docx  # python-docx for Word documents
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl  # For Excel files
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation  # python-pptx for PowerPoint
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Phishing phrases specifically for attachments - more targeted than email body
ATTACHMENT_PHISHING_KEYWORDS = {
    "credential_theft": {
        "name": "Credential Theft",
        "risk_level": "HIGH",
        "keywords": [
            "enter your password", "login required", "verify your account",
            "click to authenticate", "sign in to continue", "password expired",
            "update your credentials", "confirm your identity", "security verification required",
            "account locked", "verify identity", "re-enter password"
        ],
        "description": "Document requesting login credentials or authentication"
    },
    
    "malicious_links": {
        "name": "Malicious Links",
        "risk_level": "HIGH", 
        "keywords": [
            "click here to download", "download now", "click to verify",
            "enable content", "enable macros", "allow editing",
            "click to view document", "download attachment", "open link",
            "visit this link", "go to", "redirect to", "follow this link"
        ],
        "description": "Suspicious calls to action or link redirection"
    },
    
    "document_lures": {
        "name": "Document Lures",
        "risk_level": "HIGH",
        "keywords": [
            "document protected", "unable to display", "file corrupted",
            "enable macros to view", "content blocked", "document expired",
            "click to reload", "refresh document", "update required",
            "install plugin", "download viewer", "security warning"
        ],
        "description": "Fake document error messages designed to trick users"
    },
    
    "financial_fraud": {
        "name": "Financial Fraud",
        "risk_level": "HIGH",
        "keywords": [
            "wire transfer", "payment required", "invoice overdue", 
            "urgent payment", "bank details", "routing number",
            "account number", "swift code", "payment confirmation",
            "billing update", "payment method", "credit card"
        ],
        "description": "Financial information requests or payment redirects"
    },
    
    "executive_impersonation": {
        "name": "Executive Impersonation",
        "risk_level": "HIGH",
        "keywords": [
            "ceo request", "urgent directive", "confidential task",
            "executive order", "board decision", "management request",
            "president urgent", "chairman directive", "cfo approval",
            "immediate action required", "high priority task"
        ],
        "description": "Business Email Compromise (BEC) attempts"
    },
    
    "fake_alerts": {
        "name": "Fake Security Alerts",
        "risk_level": "MEDIUM",
        "keywords": [
            "virus detected", "security threat", "malware found",
            "computer infected", "system compromised", "scan required",
            "security update", "patch required", "vulnerability detected",
            "immediate action required", "quarantine file"
        ],
        "description": "Fake security warnings to scare users into action"
    },
    
    "social_engineering": {
        "name": "Social Engineering",
        "risk_level": "MEDIUM",
        "keywords": [
            "congratulations", "you have won", "claim your prize",
            "limited time offer", "act now", "expires soon",
            "final notice", "last chance", "don't miss out",
            "exclusive offer", "selected winner", "claim reward"
        ],
        "description": "Social engineering tactics to manipulate users"
    },
    
    "compliance_threats": {
        "name": "Compliance Threats",
        "risk_level": "MEDIUM",
        "keywords": [
            "legal action", "court notice", "compliance violation",
            "tax refund", "irs notice", "government notice",
            "fine imposed", "penalty notice", "audit required",
            "regulatory notice", "official notice", "summons"
        ],
        "description": "Fake legal or compliance threats"
    }
}

# Suspicious file extensions that can contain text/macros
ANALYZABLE_EXTENSIONS = {
    # Documents
    'pdf', 'doc', 'docx', 'docm', 'dot', 'dotx', 'dotm',
    'xls', 'xlsx', 'xlsm', 'xlt', 'xltx', 'xltm',
    'ppt', 'pptx', 'pptm', 'pot', 'potx', 'potm',
    'odt', 'ods', 'odp',  # OpenOffice/LibreOffice
    'rtf', 'txt',
    
    # Archives (may contain documents)
    'zip', 'rar', '7z',
    
    # Images (OCR capable)
    'jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'tif',
    
    # Web/Script files
    'html', 'htm', 'xml', 'js', 'vbs', 'ps1'
}

def check_content_analysis_dependencies():
    """Check which content analysis features are available."""
    available = {}
    missing = []
    
    available['pdf'] = PYMUPDF_AVAILABLE
    if not PYMUPDF_AVAILABLE:
        missing.append("PyMuPDF (pip install PyMuPDF) - for PDF analysis")
    
    available['docx'] = DOCX_AVAILABLE
    if not DOCX_AVAILABLE:
        missing.append("python-docx (pip install python-docx) - for Word document analysis")
    
    available['excel'] = OPENPYXL_AVAILABLE
    if not OPENPYXL_AVAILABLE:
        missing.append("openpyxl (pip install openpyxl) - for Excel analysis")
    
    available['powerpoint'] = PPTX_AVAILABLE
    if not PPTX_AVAILABLE:
        missing.append("python-pptx (pip install python-pptx) - for PowerPoint analysis")
    
    available['ocr'] = OCR_AVAILABLE
    if not OCR_AVAILABLE:
        missing.append("pytesseract + Pillow (pip install pytesseract Pillow) - for image OCR")
    
    # NEW: Check URL analysis availability
    available['url_analysis'] = URL_ANALYSIS_AVAILABLE
    if not URL_ANALYSIS_AVAILABLE:
        missing.append("url_extractor module - for URL reputation checking")
    
    # NEW: Check defanging availability  
    available['defanging'] = DEFANGER_AVAILABLE
    if not DEFANGER_AVAILABLE:
        missing.append("defanger module - for safe URL display")
    
    return available, missing

def extract_text_from_pdf(content):
    """Extract text from PDF using PyMuPDF."""
    if not PYMUPDF_AVAILABLE:
        return "", "PyMuPDF not available"
    
    try:
        pdf_document = fitz.open(stream=content, filetype="pdf")
        text_content = ""
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            page_text = page.get_text()
            text_content += page_text + "\n"
        
        pdf_document.close()
        return text_content.strip(), None
        
    except Exception as e:
        return "", f"PDF extraction error: {e}"

def extract_text_from_docx(content):
    """Extract text from Word document using python-docx."""
    if not DOCX_AVAILABLE:
        return "", "python-docx not available"
    
    try:
        doc = docx.Document(io.BytesIO(content))
        text_content = ""
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            text_content += paragraph.text + "\n"
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text_content += cell.text + " "
                text_content += "\n"
        
        return text_content.strip(), None
        
    except Exception as e:
        return "", f"DOCX extraction error: {e}"

def extract_text_from_excel(content):
    """Extract text from Excel file using openpyxl."""
    if not OPENPYXL_AVAILABLE:
        return "", "openpyxl not available"
    
    try:
        from openpyxl import load_workbook
        workbook = load_workbook(io.BytesIO(content), data_only=True)
        text_content = ""
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_content += f"Sheet: {sheet_name}\n"
            
            for row in sheet.iter_rows(values_only=True):
                for cell_value in row:
                    if cell_value is not None:
                        text_content += str(cell_value) + " "
                text_content += "\n"
        
        return text_content.strip(), None
        
    except Exception as e:
        return "", f"Excel extraction error: {e}"

def extract_text_from_powerpoint(content):
    """Extract text from PowerPoint using python-pptx."""
    if not PPTX_AVAILABLE:
        return "", "python-pptx not available"
    
    try:
        presentation = Presentation(io.BytesIO(content))
        text_content = ""
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text_content += f"Slide {slide_num}:\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += shape.text + "\n"
        
        return text_content.strip(), None
        
    except Exception as e:
        return "", f"PowerPoint extraction error: {e}"

def extract_text_from_image_ocr(content):
    """Extract text from image using OCR."""
    if not OCR_AVAILABLE:
        return "", "OCR libraries not available"
    
    try:
        image = Image.open(io.BytesIO(content))
        # Convert to RGB if necessary
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Extract text using OCR
        text_content = pytesseract.image_to_string(image)
        return text_content.strip(), None
        
    except Exception as e:
        return "", f"OCR extraction error: {e}"

def extract_text_from_rtf(content):
    """Basic RTF text extraction (simple approach)."""
    try:
        # Convert bytes to string
        if isinstance(content, bytes):
            rtf_text = content.decode('utf-8', errors='ignore')
        else:
            rtf_text = str(content)
        
        # Very basic RTF parsing - remove control codes
        # This is a simplified approach; a full RTF parser would be more accurate
        text_content = re.sub(r'\\[a-z]+\d*\s?', '', rtf_text)
        text_content = re.sub(r'[{}]', '', text_content)
        text_content = re.sub(r'\s+', ' ', text_content)
        
        return text_content.strip(), None
        
    except Exception as e:
        return "", f"RTF extraction error: {e}"

def extract_text_from_archive(content, filename):
    """Extract text from archive files by analyzing contained files."""
    try:
        if filename.lower().endswith('.zip'):
            with zipfile.ZipFile(io.BytesIO(content)) as zip_file:
                text_content = ""
                file_list = []
                
                for file_info in zip_file.filelist:
                    file_list.append(file_info.filename)
                    
                    # Only analyze text-like files in archives
                    if any(file_info.filename.lower().endswith(ext) for ext in ['txt', 'html', 'xml', 'js']):
                        try:
                            file_content = zip_file.read(file_info.filename)
                            file_text = file_content.decode('utf-8', errors='ignore')
                            text_content += f"File: {file_info.filename}\n{file_text}\n\n"
                        except Exception:
                            continue
                
                # Include file listing as it might contain suspicious names
                text_content = f"Archive contents: {', '.join(file_list)}\n\n" + text_content
                return text_content.strip(), None
        else:
            return "", "Archive type not supported"
            
    except Exception as e:
        return "", f"Archive extraction error: {e}"

def extract_text_from_html_xml(content):
    """Extract text from HTML/XML files."""
    try:
        if isinstance(content, bytes):
            text = content.decode('utf-8', errors='ignore')
        else:
            text = str(content)
        
        # Remove HTML/XML tags but keep the text content
        clean_text = re.sub(r'<[^>]+>', ' ', text)
        clean_text = re.sub(r'\s+', ' ', clean_text)
        
        return clean_text.strip(), None
        
    except Exception as e:
        return "", f"HTML/XML extraction error: {e}"

def extract_text_from_attachment(attachment_data):
    """
    Main function to extract text from various attachment types.
    
    Args:
        attachment_data: Dictionary containing attachment info (filename, content, etc.)
    
    Returns:
        tuple: (extracted_text, error_message)
    """
    filename = attachment_data.get('filename', '').lower()
    content = attachment_data.get('content', b'')
    
    if not content:
        return "", "No content available"
    
    # Determine extraction method based on file extension
    if filename.endswith('.pdf'):
        return extract_text_from_pdf(content)
    
    elif filename.endswith(('.doc', '.docx', '.docm')):
        if filename.endswith('.docx') or filename.endswith('.docm'):
            return extract_text_from_docx(content)
        else:
            # .doc files require different handling - could use python-docx2txt or similar
            return "", "Legacy .doc format not supported (use .docx)"
    
    elif filename.endswith(('.xls', '.xlsx', '.xlsm')):
        return extract_text_from_excel(content)
    
    elif filename.endswith(('.ppt', '.pptx', '.pptm')):
        return extract_text_from_powerpoint(content)
    
    elif filename.endswith('.rtf'):
        return extract_text_from_rtf(content)
    
    elif filename.endswith(('.txt', '.text')):
        try:
            if isinstance(content, bytes):
                return content.decode('utf-8', errors='ignore'), None
            else:
                return str(content), None
        except Exception as e:
            return "", f"Text file error: {e}"
    
    elif filename.endswith(('.html', '.htm', '.xml')):
        return extract_text_from_html_xml(content)
    
    elif filename.endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif')):
        return extract_text_from_image_ocr(content)
    
    elif filename.endswith(('.zip', '.rar', '.7z')):
        return extract_text_from_archive(content, filename)
    
    else:
        return "", f"File type not supported for text extraction: {filename}"

def extract_urls_from_text(text):
    """Extract URLs from attachment text content using existing URL patterns."""
    if not text or not isinstance(text, str):
        return []
    
    # Reuse the same URL patterns from url_extractor for consistency
    url_patterns = [
        r'https?://[^\s<>"\']+',                    # Standard HTTP URLs
        r'www\.[^\s<>"\']+',                        # www domains
        r'[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}(?:/[^\s]*)?'  # Domain patterns
    ]
    
    found_urls = []
    
    for pattern in url_patterns:
        try:
            matches = re.findall(pattern, text, re.IGNORECASE)
            found_urls.extend(matches)
        except Exception:
            continue
    
    # Clean and validate URLs (reuse logic from url_extractor)
    valid_urls = []
    for url in found_urls:
        try:
            url = url.strip()
            if len(url) > 2000:  # Skip extremely long URLs
                continue
            if url and not url.isspace():
                # Basic URL validation
                if url.startswith(('http://', 'https://', 'www.')) or '.' in url:
                    valid_urls.append(url)
        except Exception:
            continue
    
    return list(set(valid_urls))  # Remove duplicates

def analyze_urls_from_attachment(urls, api_key):
    """Analyze URLs found in attachment using existing URL analysis logic."""
    if not URL_ANALYSIS_AVAILABLE or not urls:
        return []
    
    try:
        # Group URLs by domain (reuse logic from url_extractor)
        domain_groups = defaultdict(list)
        for url in urls:
            domain = url_extractor.extract_domain(url)
            domain_groups[domain].append(url)
        
        cache = {}
        url_results = []
        
        # Analyze each domain using existing logic
        for domain, domain_urls in domain_groups.items():
            try:
                # Get representative URL (reuse logic from url_extractor)
                representative_url = url_extractor.get_shortest_url_for_domain(domain_urls)
                
                # Check VirusTotal using existing function
                verdict, comment = url_extractor.check_url_virustotal(representative_url, api_key, cache)
                
                url_results.append({
                    'domain': domain,
                    'urls': domain_urls,
                    'representative_url': representative_url,
                    'verdict': verdict,
                    'comment': comment,
                    'url_count': len(domain_urls)
                })
                
            except Exception as e:
                # Add error result
                url_results.append({
                    'domain': domain,
                    'urls': domain_urls,
                    'representative_url': domain_urls[0] if domain_urls else "",
                    'verdict': "unchecked",
                    'comment': f"Processing error: {e}",
                    'url_count': len(domain_urls)
                })
        
        return url_results
        
    except Exception as e:
        if COMPATIBLE_OUTPUT:
            print_status(f"Error analyzing URLs from attachment: {e}", "error")
        else:
            print(f"Error analyzing URLs from attachment: {e}")
        return []

def analyze_text_for_phishing(text):
    """Analyze extracted text for phishing indicators."""
    if not text or not isinstance(text, str):
        return {}
    
    findings = {}
    
    for category_id, category_data in ATTACHMENT_PHISHING_KEYWORDS.items():
        matched_keywords = []
        
        for keyword in category_data["keywords"]:
            # Use case-insensitive search
            if re.search(r'\b' + re.escape(keyword) + r'\b', text, re.IGNORECASE):
                matched_keywords.append({
                    "keyword": keyword,
                    "exact_match": True
                })
        
        if matched_keywords:
            findings[category_id] = {
                "name": category_data["name"],
                "risk_level": category_data["risk_level"],
                "description": category_data["description"],
                "matched_keywords": matched_keywords,
                "keyword_count": len(matched_keywords)
            }
    
    return findings

def calculate_attachment_risk_score(findings, text_length, url_results=None):
    """Calculate risk score for attachment content including URL analysis."""
    base_score = 0
    
    # Calculate base score from phishing keywords
    if findings:
        risk_levels = [finding["risk_level"] for finding in findings.values()]
        
        if "HIGH" in risk_levels:
            base_score = 70
            max_score = 100
        elif "MEDIUM" in risk_levels:
            base_score = 40
            max_score = 69
        else:
            max_score = 39
    
    # Factor in URL analysis results
    url_score_boost = 0
    if url_results:
        malicious_urls = sum(1 for url in url_results if url.get('verdict') == 'malicious')
        suspicious_urls = sum(1 for url in url_results if url.get('verdict') == 'suspicious')
        
        if malicious_urls > 0:
            url_score_boost = 30  # Major boost for malicious URLs
            base_score = max(base_score, 70)  # Ensure at least HIGH risk
        elif suspicious_urls > 0:
            url_score_boost = 15  # Moderate boost for suspicious URLs
            base_score = max(base_score, 40)  # Ensure at least MEDIUM risk
    
    if base_score == 0 and url_score_boost == 0:
        return 0
    
    # Calculate bonus points
    bonus = 0
    if findings:
        risk_levels = [finding["risk_level"] for finding in findings.values()]
        high_risk_count = sum(1 for level in risk_levels if level == "HIGH")
        medium_risk_count = sum(1 for level in risk_levels if level == "MEDIUM")
        total_keywords = sum(finding["keyword_count"] for finding in findings.values())
        
        if base_score >= 70:  # HIGH risk baseline
            bonus += min(20, (high_risk_count - 1) * 8)  # Additional high risk categories
            bonus += min(10, medium_risk_count * 3)      # Medium risk categories
        elif base_score >= 40:  # MEDIUM risk baseline
            bonus += min(20, (medium_risk_count - 1) * 5)  # Additional medium risk categories
        
        # Keyword density bonus (more keywords = higher risk)
        if text_length > 0:
            keyword_density = total_keywords / max(text_length / 100, 1)  # Per 100 chars
            bonus += min(10, int(keyword_density * 2))
    
    # Add URL score boost
    bonus += url_score_boost
    
    final_score = min(max_score if 'max_score' in locals() else 100, base_score + bonus)
    return final_score

def analyze_attachment_content(attachment_data, api_key=None):
    """
    Main function to analyze attachment content for phishing indicators and malicious URLs.
    
    Args:
        attachment_data: Dictionary containing attachment info
        api_key: VirusTotal API key for URL checking
    
    Returns:
        Dictionary with analysis results
    """
    filename = attachment_data.get('filename', 'unknown')
    
    # Check if file type is analyzable
    file_ext = filename.lower().split('.')[-1] if '.' in filename else ''
    if file_ext not in ANALYZABLE_EXTENSIONS:
        return {
            'analyzed': False,
            'reason': f'File type .{file_ext} not supported for content analysis',
            'text_extracted': False,
            'text_length': 0,
            'findings': {},
            'url_analysis': {},
            'risk_score': 0
        }
    
    # Extract text from attachment
    extracted_text, extraction_error = extract_text_from_attachment(attachment_data)
    
    if extraction_error:
        return {
            'analyzed': False,
            'reason': extraction_error,
            'text_extracted': False,
            'text_length': 0,
            'findings': {},
            'url_analysis': {},
            'risk_score': 0
        }
    
    if not extracted_text or len(extracted_text.strip()) < 10:
        return {
            'analyzed': True,
            'reason': 'No meaningful text content found',
            'text_extracted': True,
            'text_length': len(extracted_text) if extracted_text else 0,
            'findings': {},
            'url_analysis': {},
            'risk_score': 0
        }
    
    # Analyze extracted text for phishing indicators
    findings = analyze_text_for_phishing(extracted_text)
    
    # NEW: Extract and analyze URLs from the text
    urls = extract_urls_from_text(extracted_text)
    url_analysis = {}
    if urls:
        url_results = analyze_urls_from_attachment(urls, api_key)
        
        # Structure URL analysis similar to main URL analysis
        if url_results:
            # Group by verdict for summary
            malicious_domains = [r for r in url_results if r['verdict'] == 'malicious']
            suspicious_domains = [r for r in url_results if r['verdict'] == 'suspicious'] 
            benign_domains = [r for r in url_results if r['verdict'] == 'benign']
            unchecked_domains = [r for r in url_results if r['verdict'] == 'unchecked']
            
            url_analysis = {
                'urls_found': len(urls),
                'domains_found': len(url_results),
                'malicious_count': len(malicious_domains),
                'suspicious_count': len(suspicious_domains),
                'benign_count': len(benign_domains),
                'unchecked_count': len(unchecked_domains),
                'results': url_results
            }
    
    # Calculate risk score including URL analysis
    risk_score = calculate_attachment_risk_score(findings, len(extracted_text), url_analysis.get('results'))
    
    return {
        'analyzed': True,
        'reason': None,
        'text_extracted': True,
        'text_length': len(extracted_text),
        'extracted_text': extracted_text[:500] + '...' if len(extracted_text) > 500 else extracted_text,  # Truncate for display
        'findings': findings,
        'url_analysis': url_analysis,
        'risk_score': risk_score,
        'categories_found': len(findings)
    }

def display_attachment_content_analysis(attachment_index, filename, analysis_result):
    """Display optimized attachment content analysis results."""
    if not analysis_result.get('analyzed'):
        reason = analysis_result.get('reason', 'Unknown error')
        if COMPATIBLE_OUTPUT:
            output.print(f"- [orange3]Content Analysis: {output.escape(reason)}[/orange3]")
        else:
            print(f"- Content Analysis: {reason}")
        return
    
    findings = analysis_result.get('findings', {})
    url_analysis = analysis_result.get('url_analysis', {})
    risk_score = analysis_result['risk_score']
    text_length = analysis_result['text_length']
    
    # Always show text extraction info
    if COMPATIBLE_OUTPUT:
        output.print(f"[blue]Text extracted:[/blue] {text_length} characters")
    else:
        print(f"Text extracted: {text_length} characters")
    
    # Display URL summary if URLs found
    if url_analysis and url_analysis.get('results'):
        total_urls = url_analysis.get('urls_found', 0)
        total_domains = url_analysis.get('domains_found', 0)
        malicious_count = url_analysis.get('malicious_count', 0)
        suspicious_count = url_analysis.get('suspicious_count', 0)
        unchecked_count = url_analysis.get('unchecked_count', 0)
        benign_count = url_analysis.get('benign_count', 0)
        
        if COMPATIBLE_OUTPUT:
            output.print(f"- URLs in content: {total_urls} URL{'s' if total_urls != 1 else ''} across {total_domains} domain{'s' if total_domains != 1 else ''}")
        else:
            print(f"- URLs in content: {total_urls} URL{'s' if total_urls != 1 else ''} across {total_domains} domain{'s' if total_domains != 1 else ''}")
        
        # Show critical findings first
        if malicious_count > 0:
            if COMPATIBLE_OUTPUT:
                output.print(f"- [red]{malicious_count} malicious domain{'s' if malicious_count != 1 else ''} detected![/red]")
            else:
                print(f"- {malicious_count} malicious domain{'s' if malicious_count != 1 else ''} detected!")
        
        if suspicious_count > 0:
            if COMPATIBLE_OUTPUT:
                output.print(f"- [orange3]{suspicious_count} suspicious domain{'s' if suspicious_count != 1 else ''} detected[/orange3]")
            else:
                print(f"- {suspicious_count} suspicious domain{'s' if suspicious_count != 1 else ''} detected")
        
        if unchecked_count > 0:
            if COMPATIBLE_OUTPUT:
                output.print(f"- {unchecked_count} domain{'s' if unchecked_count != 1 else ''} unchecked by VirusTotal")
            else:
                print(f"- {unchecked_count} domain{'s' if unchecked_count != 1 else ''} unchecked by VirusTotal")
        
        if benign_count > 0:
            if COMPATIBLE_OUTPUT:
                output.print(f"- [green]{benign_count} domain{'s' if benign_count != 1 else ''} reported as benign[/green]")
            else:
                print(f"- {benign_count} domain{'s' if benign_count != 1 else ''} reported as benign")

def get_attachment_content_risk_factors(analysis_result):
    """Extract risk factors from attachment content analysis for consolidated risk display."""
    risk_factors = []
    
    findings = analysis_result.get('findings', {})
    url_analysis = analysis_result.get('url_analysis', {})
    risk_score = analysis_result.get('risk_score', 0)
    
    # Add phishing content factors
    if findings:
        high_risk_findings = [f for f in findings.values() if f['risk_level'] == 'HIGH']
        medium_risk_findings = [f for f in findings.values() if f['risk_level'] == 'MEDIUM']
        
        if high_risk_findings:
            finding_names = [f['name'] for f in high_risk_findings[:2]]
            if len(high_risk_findings) > 2:
                risk_factors.append(f"PHISHING CONTENT: {', '.join(finding_names)}, +{len(high_risk_findings) - 2} more")
            else:
                risk_factors.append(f"PHISHING CONTENT: {', '.join(finding_names)}")
        elif medium_risk_findings:
            finding_names = [f['name'] for f in medium_risk_findings[:2]]
            if len(medium_risk_findings) > 2:
                risk_factors.append(f"Suspicious content: {', '.join(finding_names)}, +{len(medium_risk_findings) - 2} more")
            else:
                risk_factors.append(f"Suspicious content: {', '.join(finding_names)}")
    
    # Add URL-based factors
    if url_analysis and url_analysis.get('results'):
        malicious_count = url_analysis.get('malicious_count', 0)
        suspicious_count = url_analysis.get('suspicious_count', 0)
        
        if malicious_count > 0:
            risk_factors.append(f"MALICIOUS URLs: {malicious_count} domain{'s' if malicious_count != 1 else ''}")
        elif suspicious_count > 0:
            risk_factors.append(f"Suspicious URLs: {suspicious_count} domain{'s' if suspicious_count != 1 else ''}")
    
    return risk_factors, risk_score

# Example usage and requirements
def get_required_packages():
    """Return list of packages needed for full functionality."""
    return [
        "PyMuPDF>=1.23.0",      # PDF analysis
        "python-docx>=0.8.11",  # Word documents
        "openpyxl>=3.1.0",      # Excel files
        "python-pptx>=0.6.21",  # PowerPoint files
        "pytesseract>=0.3.10",  # OCR for images
        "Pillow>=10.0.0"        # Image processing
    ]

if __name__ == "__main__":
    # Test availability
    available, missing = check_content_analysis_dependencies()
    
    print("Content Analysis Module Test")
    print("=" * 40)
    print(f"Available features: {list(available.keys())}")
    print(f"Missing dependencies: {missing}")
    print()
    print("Required packages for full functionality:")
    for pkg in get_required_packages():
        print(f"  pip install {pkg}")